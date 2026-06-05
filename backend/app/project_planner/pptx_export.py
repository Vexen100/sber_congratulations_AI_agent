from __future__ import annotations

import datetime as dt
import re
from io import BytesIO
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.project_planner.schemas import ConceptOption, ProjectReport

PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
BUDGET_CAVEAT = "Финансовая оценка предварительная и требует экспертной проверки."

_COLOR_DARK = RGBColor(31, 52, 45)
_COLOR_GREEN = RGBColor(25, 135, 84)
_COLOR_LIGHT_GREEN = RGBColor(226, 245, 236)
_COLOR_LIGHT = RGBColor(248, 250, 249)
_COLOR_MUTED = RGBColor(96, 108, 104)
_COLOR_WHITE = RGBColor(255, 255, 255)
_COLOR_BORDER = RGBColor(211, 229, 220)

_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)
_MARGIN_X = Inches(0.62)


def safe_text(value: object, limit: int = 180) -> str:
    text = "" if value is None else str(value)
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return "—"
    if limit <= 1 or len(text) <= limit:
        return text
    if limit <= 3:
        return text[:limit]

    cutoff = text[: limit - 1].rstrip()
    boundary = max(cutoff.rfind(" "), cutoff.rfind(","), cutoff.rfind(";"), cutoff.rfind(":"))
    if boundary >= max(18, int(limit * 0.6)):
        cutoff = cutoff[:boundary].rstrip(" ,;:")
    return cutoff.rstrip() + "…"


def format_money_rub(value: object) -> str:
    if value is None:
        return "—"
    try:
        amount = float(value)
    except (TypeError, ValueError):
        return safe_text(value, 60)
    return f"{amount:,.0f}".replace(",", " ") + " руб."


def format_date(value: object) -> str:
    if value is None:
        return "—"
    if isinstance(value, dt.datetime):
        return value.date().strftime("%d.%m.%Y")
    if isinstance(value, dt.date):
        return value.strftime("%d.%m.%Y")
    return safe_text(value, 40)


def truncate_items(items: Iterable[object], max_count: int, *, limit: int = 120) -> list[str]:
    result: list[str] = []
    for item in items:
        text = safe_text(item, limit)
        if text != "—":
            result.append(text)
        if len(result) >= max_count:
            break
    return result


def _short_bullets(items: Iterable[object], max_count: int, limit: int) -> list[str]:
    return truncate_items(items, max_count, limit=limit)


def _dedupe_key(value: object) -> str:
    text = "" if value is None else str(value)
    text = text.lower().replace("ё", "е")
    text = re.sub(r"\s+", " ", text).strip()
    return text.strip(" .,!?:;—-–")


def _unique_texts(items: Iterable[object], *, limit: int = 130) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for item in items:
        text = safe_text(item, limit)
        if text == "—":
            continue
        key = _dedupe_key(text)
        if not key or key in seen:
            continue
        seen.add(key)
        result.append(text)
    return result


def _normalize_name(value: object) -> str:
    return safe_text(value, 200).lower().replace("ё", "е").strip()


def _is_narrow_prefix_match(target: str, candidate: str) -> bool:
    if not target or not candidate.startswith(target):
        return False
    if len(candidate) == len(target):
        return True
    return candidate[len(target)] in " -—–:("


def _recommended_concept_match(report: ProjectReport) -> ConceptOption | None:
    target = _normalize_name(report.recommended_concept.concept_name)
    for concept in report.concepts:
        if _normalize_name(concept.name) == target:
            return concept
    for concept in report.concepts:
        if _is_narrow_prefix_match(target, _normalize_name(concept.name)):
            return concept
    return None


def _apply_run_font(run, *, size: int, bold: bool = False, color: RGBColor = _COLOR_DARK) -> None:
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_paragraph_text(
    paragraph,
    text: object,
    *,
    size: int,
    bold: bool = False,
    color: RGBColor = _COLOR_DARK,
    limit: int = 180,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    paragraph.text = safe_text(text, limit)
    paragraph.alignment = align
    for run in paragraph.runs:
        _apply_run_font(run, size=size, bold=bold, color=color)


def _text_box(
    slide,
    text: object,
    *,
    x,
    y,
    w,
    h,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = _COLOR_DARK,
    limit: int = 180,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    _set_paragraph_text(
        frame.paragraphs[0],
        text,
        size=size,
        bold=bold,
        color=color,
        limit=limit,
        align=align,
    )


def _shape_fill(shape, color: RGBColor, line_color: RGBColor | None = None) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line_color or color


def _add_section_label(slide, text: str, *, x=_MARGIN_X, y=Inches(0.24), w=Inches(2.15)) -> None:
    label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.33))
    _shape_fill(label, _COLOR_LIGHT_GREEN, _COLOR_LIGHT_GREEN)
    frame = label.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_paragraph_text(frame.paragraphs[0], text, size=10, bold=True, color=_COLOR_GREEN, limit=45)


def _add_pitch_title(slide, title: str, subtitle: str | None = None) -> None:
    _text_box(
        slide,
        title,
        x=_MARGIN_X,
        y=Inches(0.6),
        w=Inches(12.1),
        h=Inches(0.58),
        size=24,
        bold=True,
        limit=72,
    )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, _MARGIN_X, Inches(1.22), Inches(12.1), Inches(0.03)
    )
    _shape_fill(line, _COLOR_GREEN, _COLOR_GREEN)
    if subtitle:
        _text_box(
            slide,
            subtitle,
            x=_MARGIN_X,
            y=Inches(1.34),
            w=Inches(11.8),
            h=Inches(0.34),
            size=11,
            color=_COLOR_MUTED,
            limit=125,
        )


def _add_card(
    slide,
    title: str,
    lines: Sequence[object],
    *,
    x,
    y,
    w,
    h,
    accent: bool = False,
    title_size: int = 13,
    body_size: int = 11,
    max_lines: int = 3,
    body_limit: int = 110,
) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _shape_fill(card, _COLOR_LIGHT_GREEN if accent else _COLOR_WHITE, _COLOR_BORDER)
    frame = card.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.17)
    frame.margin_right = Inches(0.15)
    frame.margin_top = Inches(0.11)
    frame.margin_bottom = Inches(0.08)
    _set_paragraph_text(
        frame.paragraphs[0],
        title,
        size=title_size,
        bold=True,
        color=_COLOR_GREEN if accent else _COLOR_DARK,
        limit=62,
    )
    for line in _short_bullets(lines, max_lines, body_limit):
        paragraph = frame.add_paragraph()
        paragraph.space_before = Pt(4)
        _set_paragraph_text(paragraph, line, size=body_size, color=_COLOR_DARK, limit=body_limit)


def _add_metric_card(slide, label: str, value: object, *, x, y, w, h) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _shape_fill(card, _COLOR_WHITE, _COLOR_BORDER)
    frame = card.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.16)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.1)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    _set_paragraph_text(
        frame.paragraphs[0], label, size=10, bold=True, color=_COLOR_GREEN, limit=35
    )
    value_paragraph = frame.add_paragraph()
    _set_paragraph_text(value_paragraph, value, size=15, bold=True, limit=48)


def _add_value_card(
    slide, title: str, lines: Sequence[object], *, x, y, w, h, accent=False
) -> None:
    _add_card(
        slide,
        title,
        lines,
        x=x,
        y=y,
        w=w,
        h=h,
        accent=accent,
        title_size=14 if accent else 12,
        body_size=12 if accent else 11,
        max_lines=3,
        body_limit=118,
    )


def _new_slide(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _COLOR_LIGHT
    _add_pitch_title(slide, title, subtitle)
    return slide


def _cover_type(report: ProjectReport) -> str:
    text = " ".join(
        [
            report.source_input.idea,
            report.passport.title,
            report.source_input.project_accents or "",
        ]
    ).lower()
    if any(marker in text for marker in ("фестиваль", "мероприят", "форум")):
        return "Мероприятие"
    if any(marker in text for marker in ("сервис", "портал", "интеграц")):
        return "IT / сервис"
    return "Проектная инициатива"


def _project_value_statement(report: ProjectReport) -> str:
    relevance = safe_text(report.passport.relevance_for_ural_bank, 115)
    if relevance != "—":
        return relevance
    return safe_text(report.passport.goal, 115)


def _slide_title(prs: Presentation, report: ProjectReport) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _COLOR_LIGHT
    _add_section_label(slide, "Проектная инициатива", y=Inches(0.42), w=Inches(2.45))

    _text_box(
        slide,
        report.passport.title,
        x=Inches(0.76),
        y=Inches(0.96),
        w=Inches(11.85),
        h=Inches(1.1),
        size=31,
        bold=True,
        limit=86,
    )
    _text_box(
        slide,
        "Презентация для защиты проекта",
        x=Inches(0.78),
        y=Inches(2.02),
        w=Inches(7.0),
        h=Inches(0.35),
        size=15,
        color=_COLOR_GREEN,
        limit=80,
    )

    metrics = [
        ("География", report.source_input.geography),
        ("Срок", format_date(report.source_input.deadline)),
        ("Бюджет", format_money_rub(report.resources.financial_total)),
        ("Формат", _cover_type(report)),
    ]
    x = Inches(0.76)
    for label, value in metrics:
        _add_metric_card(slide, label, value, x=x, y=Inches(2.85), w=Inches(2.9), h=Inches(1.05))
        x += Inches(3.02)

    _add_value_card(
        slide,
        "Главная мысль",
        [_project_value_statement(report)],
        x=Inches(0.76),
        y=Inches(4.55),
        w=Inches(11.85),
        h=Inches(1.25),
        accent=True,
    )


def _slide_need(prs: Presentation, report: ProjectReport) -> None:
    slide = _new_slide(prs, "Зачем проект нужен", "Фокус на управленческой ценности")
    success_criteria = _unique_texts(report.passport.success_criteria, limit=115)
    expected_effect = success_criteria[:2]
    success_checks = success_criteria[2:5]
    if not success_checks:
        success_checks = ["Детальные критерии — в DOCX."]
    _add_value_card(
        slide,
        "Возможность / задача",
        [report.passport.goal, report.source_input.project_accents],
        x=Inches(0.72),
        y=Inches(1.78),
        w=Inches(11.85),
        h=Inches(1.65),
        accent=True,
    )
    _add_value_card(
        slide,
        "Ценность для банка",
        [report.passport.relevance_for_ural_bank],
        x=Inches(0.72),
        y=Inches(4.0),
        w=Inches(3.75),
        h=Inches(1.6),
    )
    _add_value_card(
        slide,
        "Ожидаемый эффект",
        expected_effect,
        x=Inches(4.78),
        y=Inches(4.0),
        w=Inches(3.75),
        h=Inches(1.6),
    )
    _add_value_card(
        slide,
        "Критерии успеха",
        success_checks,
        x=Inches(8.84),
        y=Inches(4.0),
        w=Inches(3.75),
        h=Inches(1.6),
    )


def _slide_solution(prs: Presentation, report: ProjectReport) -> None:
    concept = _recommended_concept_match(report)
    slide = _new_slide(prs, "Предлагаемое решение")
    _add_section_label(
        slide, "Рекомендуем к реализации", x=Inches(0.72), y=Inches(1.45), w=Inches(2.75)
    )
    if concept is None:
        _add_value_card(
            slide,
            "Концепция требует уточнения",
            ["Рекомендуемая концепция требует уточнения в полном отчёте."],
            x=Inches(0.72),
            y=Inches(1.95),
            w=Inches(11.85),
            h=Inches(1.8),
            accent=True,
        )
        return

    _add_value_card(
        slide,
        concept.name,
        [concept.key_idea, f"Почему: {safe_text(report.recommended_concept.rationale, 118)}"],
        x=Inches(0.72),
        y=Inches(1.95),
        w=Inches(11.85),
        h=Inches(2.0),
        accent=True,
    )
    for index, advantage in enumerate(_short_bullets(concept.advantages, 3, 82)):
        _add_card(
            slide,
            f"Польза {index + 1}",
            [advantage],
            x=Inches(0.72 + index * 4.06),
            y=Inches(4.55),
            w=Inches(3.74),
            h=Inches(1.25),
            title_size=12,
            body_size=11,
            max_lines=1,
            body_limit=82,
        )


def _add_concept_choice_cards(slide, report: ProjectReport) -> None:
    recommended = _recommended_concept_match(report)
    visible = list(report.concepts[:3])
    if recommended and recommended not in visible:
        visible = [recommended, *visible[:2]]
    if not visible:
        _add_card(
            slide,
            "Концепции",
            ["Концепции требуют уточнения в полном отчёте."],
            x=Inches(0.75),
            y=Inches(1.85),
            w=Inches(11.75),
            h=Inches(1.55),
        )
        return

    for index, concept in enumerate(visible[:3]):
        is_recommended = recommended is concept
        lines = [
            f"Бюджет: {format_money_rub(concept.estimated_cost)} · Сложность: {safe_text(concept.effort_level, 30)}",
            "Сильная сторона: "
            + safe_text(concept.advantages[0] if concept.advantages else None, 66),
            "Компромисс: "
            + safe_text(concept.disadvantages[0] if concept.disadvantages else None, 66),
        ]
        title = (
            f"Рекомендуемый вариант · {safe_text(concept.name, 45)}"
            if is_recommended
            else concept.name
        )
        _add_card(
            slide,
            title,
            lines,
            x=Inches(0.72 + index * 4.06),
            y=Inches(1.85),
            w=Inches(3.74),
            h=Inches(3.55),
            accent=is_recommended,
            title_size=12,
            body_size=11,
            max_lines=3,
            body_limit=78,
        )


def _slide_choice(prs: Presentation, report: ProjectReport) -> None:
    slide = _new_slide(prs, "Почему этот вариант", "Сравнение только на уровне решения")
    _add_concept_choice_cards(slide, report)
    if len(report.concepts) > 3:
        _text_box(
            slide,
            "Полное сравнение — в DOCX.",
            x=Inches(0.78),
            y=Inches(6.18),
            w=Inches(5.0),
            h=Inches(0.3),
            size=11,
            color=_COLOR_MUTED,
            limit=60,
        )


def _add_timeline_cards(slide, report: ProjectReport) -> None:
    phases = report.roadmap[:4]
    if not phases:
        _add_card(
            slide,
            "План реализации",
            ["Фазы требуют уточнения в полном отчёте."],
            x=Inches(0.75),
            y=Inches(1.85),
            w=Inches(11.75),
            h=Inches(1.5),
        )
        return

    for index, phase in enumerate(phases):
        x = Inches(0.72 + index * 3.0)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1.23), Inches(1.77), Inches(0.2), Inches(0.2)
        )
        _shape_fill(dot, _COLOR_GREEN, _COLOR_GREEN)
        if index < len(phases) - 1:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x + Inches(1.45), Inches(1.86), Inches(2.62), Inches(0.03)
            )
            _shape_fill(line, _COLOR_BORDER, _COLOR_BORDER)
        milestone = (
            phase.milestones[0].title if phase.milestones else "Контрольная точка требует уточнения"
        )
        lines = [
            f"{format_date(phase.start_date)} — {format_date(phase.end_date)}",
            "Ключевая точка: " + safe_text(milestone, 74),
        ]
        _add_card(
            slide,
            phase.name,
            lines,
            x=x,
            y=Inches(2.1),
            w=Inches(2.84),
            h=Inches(3.15),
            accent=index == 0,
            title_size=13,
            body_size=11,
            max_lines=2,
            body_limit=82,
        )
    if len(report.roadmap) > 4 or any(len(phase.milestones) > 1 for phase in report.roadmap):
        _text_box(
            slide,
            "Полный план — в DOCX.",
            x=Inches(0.78),
            y=Inches(6.0),
            w=Inches(5.0),
            h=Inches(0.3),
            size=11,
            color=_COLOR_MUTED,
            limit=60,
        )


def _slide_execution(prs: Presentation, report: ProjectReport) -> None:
    slide = _new_slide(prs, "Как реализуем", "План достаточно конкретен для запуска")
    _add_timeline_cards(slide, report)


def _add_budget_pitch(slide, report: ProjectReport) -> None:
    _add_metric_card(
        slide,
        "Бюджетная рамка",
        format_money_rub(report.resources.financial_total),
        x=Inches(0.75),
        y=Inches(1.75),
        w=Inches(3.85),
        h=Inches(1.3),
    )
    for index, item in enumerate(report.resources.financial_items[:4]):
        _add_card(
            slide,
            safe_text(item.category, 52),
            [format_money_rub(item.amount)],
            x=Inches(0.75),
            y=Inches(3.35 + index * 0.66),
            w=Inches(3.85),
            h=Inches(0.54),
            title_size=9,
            body_size=10.8,
            max_lines=1,
            body_limit=45,
        )
    if len(report.resources.financial_items) > 4:
        _text_box(
            slide,
            "Остальные статьи — в DOCX.",
            x=Inches(0.78),
            y=Inches(6.1),
            w=Inches(4.0),
            h=Inches(0.3),
            size=11,
            color=_COLOR_MUTED,
            limit=60,
        )
    _add_value_card(
        slide,
        "Материальные ресурсы",
        report.resources.material_resources,
        x=Inches(5.0),
        y=Inches(1.75),
        w=Inches(3.55),
        h=Inches(1.85),
    )
    _add_value_card(
        slide,
        "Информационные ресурсы",
        report.resources.information_resources,
        x=Inches(8.95),
        y=Inches(1.75),
        w=Inches(3.55),
        h=Inches(1.85),
    )
    _add_value_card(
        slide,
        "Оговорка по оценке",
        [BUDGET_CAVEAT],
        x=Inches(5.0),
        y=Inches(4.3),
        w=Inches(7.5),
        h=Inches(1.35),
        accent=True,
    )


def _slide_investment(prs: Presentation, report: ProjectReport) -> None:
    slide = _new_slide(prs, "Инвестиции и ресурсы")
    _add_budget_pitch(slide, report)


def _add_decision_slide(slide, report: ProjectReport) -> None:
    _add_section_label(slide, "Ключевые роли", x=Inches(0.72), y=Inches(1.42), w=Inches(1.8))
    for index, role in enumerate(report.team[:3]):
        _add_card(
            slide,
            role.title,
            [f"{role.count} чел.", safe_text(", ".join(role.competencies[:2]), 64)],
            x=Inches(0.72),
            y=Inches(1.86 + index * 0.82),
            w=Inches(4.75),
            h=Inches(0.65),
            title_size=10,
            body_size=9.5,
            max_lines=2,
            body_limit=70,
        )

    _add_section_label(
        slide, "Первые управленческие шаги", x=Inches(6.0), y=Inches(1.42), w=Inches(3.0)
    )
    first_phase = report.roadmap[0] if report.roadmap else None
    first_milestone = (
        first_phase.milestones[0].title if first_phase and first_phase.milestones else None
    )
    decision_lines = [
        "Утвердить рекомендуемую концепцию.",
        "Назначить ответственных за стартовую фазу.",
        "Зафиксировать бюджетную рамку и контрольную точку.",
    ]
    if first_milestone:
        decision_lines[2] = "Контрольная точка: " + safe_text(first_milestone, 70)
    _add_value_card(
        slide,
        "Что нужно решить",
        decision_lines,
        x=Inches(6.0),
        y=Inches(1.86),
        w=Inches(6.35),
        h=Inches(1.65),
        accent=True,
    )

    for index, item in enumerate(report.raci[:2]):
        _add_card(
            slide,
            safe_text(item.activity, 60),
            [
                f"Активность → R: {safe_text(item.responsible, 36)} / A: {safe_text(item.accountable, 36)}"
            ],
            x=Inches(6.0),
            y=Inches(3.85 + index * 0.7),
            w=Inches(6.35),
            h=Inches(0.56),
            title_size=9.5,
            body_size=9.5,
            max_lines=1,
            body_limit=98,
        )
    if len(report.team) > 3 or len(report.raci) > 2:
        _text_box(
            slide,
            "Полная матрица RACI — в DOCX.",
            x=Inches(6.05),
            y=Inches(5.15),
            w=Inches(5.5),
            h=Inches(0.28),
            size=11,
            color=_COLOR_MUTED,
            limit=70,
        )
    _add_value_card(
        slide,
        "Следующий шаг",
        ["Утвердить концепцию, ответственных и бюджетную рамку."],
        x=Inches(0.72),
        y=Inches(5.78),
        w=Inches(11.65),
        h=Inches(0.95),
        accent=True,
    )


def _slide_decision(prs: Presentation, report: ProjectReport) -> None:
    slide = _new_slide(prs, "Что нужно для запуска")
    _add_decision_slide(slide, report)


def export_project_report_pptx(report: ProjectReport) -> bytes:
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    prs.core_properties.author = "Project Planner"
    prs.core_properties.title = safe_text(report.passport.title, 120)
    prs.core_properties.subject = "Project Planner PPTX export"
    fixed_timestamp = dt.datetime(2026, 1, 1, 0, 0, 0)
    prs.core_properties.created = fixed_timestamp
    prs.core_properties.modified = fixed_timestamp

    _slide_title(prs, report)
    _slide_need(prs, report)
    _slide_solution(prs, report)
    _slide_choice(prs, report)
    _slide_execution(prs, report)
    _slide_investment(prs, report)
    if report.team or report.raci:
        _slide_decision(prs, report)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
